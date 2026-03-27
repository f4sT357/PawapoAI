import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import List, Optional
from utils.logger import logger

class PPTXService:
    @staticmethod
    def create_pptx_with_template(slides_data: List[dict], template_path: Optional[str] = None) -> bytes:
        """
        スライド詳細データからPowerPointを生成する
        """
        try:
            # 1. テンプレートの読み込み
            if template_path:
                prs = Presentation(template_path)
                start_slide_idx = 1  # 2ページ目から開始（1ページ目は表紙）
            else:
                prs = Presentation()
                start_slide_idx = 0

            # 2. スライドの追加/更新
            for idx, slide_info in enumerate(slides_data):
                slide_idx = start_slide_idx + idx
                
                # 既存スライドの数を超えた場合は新規スライドを追加
                if slide_idx >= len(prs.slides):
                    try:
                        # テンプレートのレイアウト 1 (通常はタイトルと本文) を使用
                        slide_layout = prs.slide_layouts[1]
                    except IndexError:
                        slide_layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(slide_layout)
                else:
                    # 既存スライドを上書き
                    slide = prs.slides[slide_idx]
                
                # 3. タイトルの入力
                title_placeholder = slide.shapes.title
                if title_placeholder:
                    title_placeholder.text = slide_info.get("title", "")
                    # タイトルの位置調整 (必要に応じて)
                    title_placeholder.top = Inches(0.2)
                    title_placeholder.left = Inches(0.5)
                    title_placeholder.width = Inches(9)
                    title_placeholder.height = Inches(0.8)
                
                # 4. 本文の入力
                body_placeholder = None
                # タイトル以外の、テキストを書き込めるプレースホルダーを探す
                for shape in slide.placeholders:
                    if shape.has_text_frame and shape.placeholder_format.idx != 0:
                        body_placeholder = shape
                        break
                
                if body_placeholder:
                    # 本文の位置調整
                    body_placeholder.top = Inches(2.0)
                    body_placeholder.left = Inches(0.5)
                    body_placeholder.width = Inches(9)
                    body_placeholder.height = Inches(4.5)
                    
                    tf = body_placeholder.text_frame
                    tf.clear()
                    tf.word_wrap = True
                    
                    for point in slide_info.get("content", []):
                        p = tf.add_paragraph()
                        p.text = str(point)
                        p.level = 0
                        p.alignment = PP_ALIGN.LEFT
                        
                        # テキストカラーとサイズの設定
                        for run in p.runs:
                            run.font.color.rgb = RGBColor(0, 0, 0)  # ブラック
                            run.font.size = Pt(18)  # 文字サイズ18ポイント
                
            # 5. バイトデータとして保存
            binary_io = io.BytesIO()
            prs.save(binary_io)
            return binary_io.getvalue()

        except Exception as e:
            logger.error(f"PPTX generation error: {e}")
            raise
